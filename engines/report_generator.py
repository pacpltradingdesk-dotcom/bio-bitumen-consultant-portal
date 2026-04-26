"""
Report Generator — Auto-generate PPT + DPR PDF from cfg
=========================================================
Produces bank-ready, investor-ready documents with:
- Write-up + Numbers + Graphs + Images
- All values from cfg (single source of truth)
- Versioned exports with audit trail
"""
import io
import os
from pathlib import Path
from datetime import datetime

EXPORT_DIR = Path(__file__).parent.parent / "data" / "exports"


def _ensure():
    EXPORT_DIR.mkdir(parents=True, exist_ok=True)


# ══════════════════════════════════════════════════════════════════════
# PPT AUTO-GENERATOR (25 slides from cfg)
# ══════════════════════════════════════════════════════════════════════
def generate_pptx(cfg, company=None):
    """Generate a complete .pptx presentation from cfg.
    Returns: (BytesIO buffer, filename) or (None, error)
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        return None, "python-pptx not installed. Run: pip install python-pptx"

    if company is None:
        try:
            from config import COMPANY
            company = COMPANY
        except Exception:
            company = {"trade_name": "PPS Anantams", "owner": "Prince P. Shah", "phone": ""}

    tpd = cfg.get("capacity_tpd", 20)
    state = cfg.get("state", "Maharashtra")
    inv = cfg.get("investment_cr", 8)
    roi = cfg.get("roi_pct", 20)
    irr = cfg.get("irr_pct", 26)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    def add_slide(title, subtitle="", body=""):
        layout = prs.slide_layouts[1]  # Title + Content
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        if subtitle:
            slide.placeholders[1].text = subtitle + "\n\n" + body
        return slide

    # Slide 1: Cover
    add_slide(
        f"Bio-Bitumen Plant Project — {tpd:.0f} TPD",
        f"Prepared for: {cfg.get('client_name', 'Valued Client')}",
        f"By: {company.get('trade_name', 'PPS Anantams')}\n"
        f"Date: {datetime.now().strftime('%d %B %Y')}\n"
        f"Location: {cfg.get('location', '')}, {state}"
    )

    # Slide 2: Project Overview
    add_slide("Project Overview",
        f"Technology: Biomass Pyrolysis (CSIR-CRRI Licensed)\n"
        f"Capacity: {tpd:.0f} TPD agro-waste feed\n"
        f"Investment: Rs {inv:.2f} Crore\n"
        f"Location: {cfg.get('location', '')}, {state}\n"
        f"Products: Bio-Bitumen VG30/VG40, Bio-Char, Bio-Oil, Carbon Credits"
    )

    # Slide 3: Market Opportunity
    add_slide("Market Opportunity",
        "India Bitumen Market: Rs 25,000+ Crore/year\n"
        "Import Dependency: 49%\n"
        "Plants Needed: 130-216 in next 5-7 years\n"
        "NHAI Green Mandate: 15% bio-bitumen by 2030\n"
        "CSIR-CRRI technology transferred to 14 companies (Jan 2026)"
    )

    # Slide 4: Technology
    add_slide("Technology — CSIR-CRRI Licensed",
        f"Process: Biomass → Pyrolysis (500°C) → Bio-Oil + Char + Syngas\n"
        f"Yields: Bio-Oil {cfg.get('bio_oil_yield_pct', 32)}% | "
        f"Bio-Char {cfg.get('bio_char_yield_pct', 28)}% | "
        f"Syngas {cfg.get('syngas_yield_pct', 22)}%\n"
        f"Blend: {cfg.get('bio_blend_pct', 20)}% bio-oil in VG30 bitumen\n"
        f"Quality: Meets IS:73 VG30/VG40 specifications"
    )

    # Slide 5: Process Flow
    add_slide("13-Step Manufacturing Process",
        "1. Biomass Receiving → 2. Shredding → 3. Drying\n"
        "4. Pelletization → 5. Pyrolysis (450-550°C)\n"
        "6. Bio-Oil Condensation → 7. Oil Storage\n"
        "8. VG-30 Heating → 9. High Shear Blending\n"
        "10. Quality Testing (IS:73) → 11. Product Storage\n"
        "12. Drum/Tanker Packing → 13. Dispatch to NHAI/PWD"
    )

    # Slide 6-7: Equipment
    try:
        from engines.plant_engineering import compute_all, get_machinery_list
        comp = compute_all(cfg)
        mach = get_machinery_list(cfg, comp)
        mach_text = "\n".join(f"{m['tag']}: {m['name']} — {m['dims'][:50]}" for m in mach[:8])
        add_slide(f"Key Equipment — {len(mach)} Items",
            f"Reactor: {comp['reactor_qty']}x Ø{comp['reactor_dia_m']}m × {comp['reactor_ht_m']}m\n"
            f"Dryer: Ø{comp['dryer_dia_m']}m × {comp['dryer_len_m']}m\n"
            f"Total Motor Load: {sum(m.get('motor_kw',0)*m.get('qty',1) for m in mach)} kW\n\n"
            f"{mach_text}"
        )
    except Exception:
        add_slide("Equipment", f"Capacity: {tpd:.0f} TPD — equipment list based on capacity")

    # Slide 8: BOQ
    try:
        from state_manager import calculate_boq
        boq = calculate_boq(tpd, process_id=cfg.get("process_id", 1))
        total_lac = sum(i["amount_lac"] for i in boq)
        cats = {}
        for i in boq:
            cats[i["category"]] = cats.get(i["category"], 0) + i["amount_lac"]
        boq_text = "\n".join(f"{c}: Rs {v:.0f} Lac" for c, v in sorted(cats.items())[:8])
        add_slide(f"Bill of Quantities — {len(boq)} Items",
            f"Total BOQ: Rs {total_lac:.0f} Lac ({total_lac/100:.2f} Cr)\n"
            f"Zones: {len(cats)}\n\n{boq_text}"
        )
    except Exception:
        add_slide("BOQ", f"Investment: Rs {inv:.2f} Crore")

    # Slide 9: Plant Layout
    add_slide(f"Plant Layout — {cfg.get('plot_length_m', 120)}m × {cfg.get('plot_width_m', 80)}m",
        f"15 Zones: Gate → RM → Processing → Reactor → Oil → Blending → "
        f"Storage → Dispatch → Electrical → Utilities → Lab → Safety → "
        f"Civil → Office → Maintenance\n\n"
        f"Safety: 15m reactor clearance, 6m roads, 45m hydrant spacing"
    )

    # Slide 10-14: Financials
    add_slide("Project Cost",
        f"Total Investment: Rs {inv:.2f} Crore\n"
        f"Bank Loan ({(1-cfg.get('equity_ratio', 0.4))*100:.0f}%): Rs {cfg.get('loan_cr', inv*0.6):.2f} Crore\n"
        f"Promoter Equity ({cfg.get('equity_ratio', 0.4)*100:.0f}%): Rs {cfg.get('equity_cr', inv*0.4):.2f} Crore\n"
        f"EMI: Rs {cfg.get('emi_lac_mth', 0):.2f} Lac/month\n"
        f"Interest Rate: {cfg.get('interest_rate', 0.115)*100:.1f}% p.a."
    )

    add_slide("Revenue & Profitability",
        f"Revenue Year 5: Rs {cfg.get('revenue_yr5_lac', 0):.0f} Lac\n"
        f"Monthly Profit: Rs {cfg.get('monthly_profit_lac', 0):.1f} Lac\n"
        f"ROI: {roi:.1f}%\n"
        f"IRR: {irr:.1f}%\n"
        f"Break-Even: {cfg.get('break_even_months', 0)} months\n"
        f"DSCR Year 3: {cfg.get('dscr_yr3', 0):.2f}x"
    )

    # Slide 15: Compliance
    add_slide("Compliance & Licenses",
        "Key Approvals Required:\n"
        "- CTE (Consent to Establish) — State PCB\n"
        "- CTO (Consent to Operate) — State PCB\n"
        "- Factory License — Factory Inspector\n"
        "- Fire NOC — Fire Department\n"
        "- PESO License — for petroleum storage\n"
        "- GST Registration\n"
        "- MSME Registration\n"
        "Total: 25 licenses — We handle ALL"
    )

    # Slide 16: Government Support
    add_slide("Government Schemes & Subsidies",
        f"MNRE Biomass: 25% capital subsidy (max Rs 2.1 Cr)\n"
        f"State MSME: Varies by state\n"
        f"CGTMSE: Collateral-free up to Rs 5 Cr\n"
        f"Carbon Credits: Rs {tpd*300*0.35*12500/100000:.1f} Lac/year\n"
        f"NHAI Green Mandate: Priority procurement"
    )

    # Slide 17: Timeline
    add_slide("Project Timeline — 12 to 18 Months",
        "Month 1: DPR + Feasibility\n"
        "Month 1-2: Company Setup\n"
        "Month 2-4: Land + Approvals\n"
        "Month 3-7: Environmental Clearances\n"
        "Month 3-5: Bank Loan\n"
        "Month 5-8: Equipment Order\n"
        "Month 6-10: Civil Construction\n"
        "Month 10-12: Installation\n"
        "Month 13+: Commercial Production"
    )

    # Slide 18: Why PPS Anantams
    add_slide(f"Why {company.get('trade_name', 'PPS Anantams')}?",
        f"Experience: {company.get('years_experience', 25)} years in bitumen industry (since 2001)\n"
        f"Director: {company.get('years_as_director', 17)}+ years MCA-registered (DIN 06680837, since 2009)\n"
        f"Group: Founder of BSE-Listed Omnipotent Industries Ltd (2016) — 1.2 Lakh MT traded, 11 JVs\n"
        f"Plants Engaged: {company.get('plants_engaged', 9)} — {company.get('plants_breakdown', '3 as GM, 1 as CEO, 3 as Founder/MD, 2 as Consultant')}\n"
        f"Network: {company.get('industry_contacts', 4452):,}-contact industry database (via Omnipotent)\n"
        f"International: VG-30 import capacity up to 2.4 Lakh MT/yr (Iraq/USA)\n"
        f"Turnkey: DPR to Dispatch — complete solution"
    )

    # Slide 19: Next Steps
    add_slide("Next Steps",
        "1. Finalize capacity and location\n"
        "2. Sign consulting agreement\n"
        "3. DPR generation (7 days)\n"
        "4. Bank loan application\n"
        "5. Site identification\n"
        "6. Equipment procurement\n"
        "7. Construction & commissioning\n\n"
        f"Contact: {company.get('owner', '')} | {company.get('phone', '')}"
    )

    # Save to buffer
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    filename = f"Bio_Bitumen_{int(tpd)}TPD_{state}_{datetime.now().strftime('%Y%m%d')}.pptx"

    # Register export
    try:
        from engines.audit_logger import register_export, log_event
        _ensure()
        fpath = EXPORT_DIR / filename
        fpath.write_bytes(buf.getvalue())
        buf.seek(0)
        register_export("pptx", str(fpath), cfg)
        log_event("EXPORT_CREATED", "report_generator", f"PPT {filename}", "success")
    except Exception:
        pass

    return buf, filename


# ══════════════════════════════════════════════════════════════════════
# DPR PDF AUTO-GENERATOR (Bank-ready from cfg)
# ══════════════════════════════════════════════════════════════════════
def generate_dpr_pdf(cfg, company=None):
    """Generate a DPR PDF document from cfg.
    Uses basic HTML→PDF approach. Returns (BytesIO, filename) or (None, error).
    """
    if company is None:
        try:
            from config import COMPANY
            company = COMPANY
        except Exception:
            company = {"trade_name": "PPS Anantams", "owner": "Prince P. Shah"}

    tpd = cfg.get("capacity_tpd", 20)
    state = cfg.get("state", "Maharashtra")
    inv = cfg.get("investment_cr", 8)

    # Build DPR HTML content
    html = f"""
    <html><head><style>
    body {{ font-family: Calibri, Arial; margin: 40px; color: #333; font-size: 12px; }}
    h1 {{ color: #003366; border-bottom: 2px solid #003366; padding-bottom: 10px; }}
    h2 {{ color: #006699; margin-top: 25px; }}
    table {{ border-collapse: collapse; width: 100%; margin: 10px 0; }}
    th {{ background: #003366; color: white; padding: 8px; text-align: left; }}
    td {{ padding: 6px 8px; border-bottom: 1px solid #ddd; }}
    .header {{ text-align: center; margin-bottom: 30px; }}
    .metric {{ display: inline-block; margin: 10px; padding: 15px; border: 1px solid #ddd; border-radius: 8px; }}
    </style></head><body>

    <div class="header">
        <h1>DETAILED PROJECT REPORT (DPR)</h1>
        <h2>{cfg.get('project_name', f'Bio-Bitumen Plant — {tpd:.0f} TPD')}</h2>
        <p>Prepared by: {company.get('trade_name', 'PPS Anantams')} | {company.get('owner', '')}</p>
        <p>Date: {datetime.now().strftime('%d %B %Y')} | Version: {cfg.get('dpr_version', 'v1.0')}</p>
    </div>

    <h2>1. Executive Summary</h2>
    <p>This DPR presents a {tpd:.0f} TPD bio-modified bitumen plant project
    in {cfg.get('location', '')}, {state} with a total investment of Rs {inv:.2f} Crore.
    The plant uses CSIR-CRRI licensed pyrolysis technology to convert agricultural
    waste into bio-bitumen conforming to IS:73 VG30/VG40 specifications.</p>

    <table>
    <tr><th>Parameter</th><th>Value</th></tr>
    <tr><td>Plant Capacity</td><td>{tpd:.0f} MT/Day (agro-waste feed)</td></tr>
    <tr><td>Total Investment</td><td>Rs {inv:.2f} Crore</td></tr>
    <tr><td>Bank Loan</td><td>Rs {cfg.get('loan_cr', inv*0.6):.2f} Crore</td></tr>
    <tr><td>Promoter Equity</td><td>Rs {cfg.get('equity_cr', inv*0.4):.2f} Crore</td></tr>
    <tr><td>ROI</td><td>{cfg.get('roi_pct', 20):.1f}%</td></tr>
    <tr><td>IRR</td><td>{cfg.get('irr_pct', 26):.1f}%</td></tr>
    <tr><td>DSCR (Year 3)</td><td>{cfg.get('dscr_yr3', 0):.2f}x</td></tr>
    <tr><td>Break-Even</td><td>{cfg.get('break_even_months', 0)} months</td></tr>
    <tr><td>Revenue Year 5</td><td>Rs {cfg.get('revenue_yr5_lac', 0):.0f} Lakhs</td></tr>
    </table>

    <h2>2. Promoter Profile</h2>
    <p>{company.get('trade_name', 'PPS Anantams Corporation Pvt Ltd')}<br>
    Owner: {company.get('owner', 'Prince Pratap Shah')}<br>
    Experience: {company.get('years_experience', 25)} years in bitumen industry (since 2001) | Director since 2009 (DIN 06680837) | Founder Omnipotent Industries (BSE-Listed, 2016)<br>
    Contact: {company.get('phone', '')} | {company.get('email', '')}</p>

    <h2>3. Technology</h2>
    <p>Pyrolysis of agricultural biomass at 450-550 degrees C in oxygen-free reactor.<br>
    Yields: Bio-Oil {cfg.get('bio_oil_yield_pct', 32)}% | Bio-Char {cfg.get('bio_char_yield_pct', 28)}% |
    Syngas {cfg.get('syngas_yield_pct', 22)}% | Loss {cfg.get('process_loss_pct', 18)}%<br>
    Blending: {cfg.get('bio_blend_pct', 20)}% bio-oil with VG30 conventional bitumen<br>
    Standard: CSIR-CRRI licensed, meets IS:73 and MoRTH 2026 Section 519.</p>

    <h2>4. Location</h2>
    <p>State: {state}<br>
    City: {cfg.get('location', '')}<br>
    Site: {cfg.get('site_address', 'To be finalized')}<br>
    Plot: {cfg.get('plot_length_m', 120)}m x {cfg.get('plot_width_m', 80)}m</p>

    <h2>5. Financial Summary</h2>
    <table>
    <tr><th>Item</th><th>Amount</th></tr>
    <tr><td>Total Investment</td><td>Rs {inv:.2f} Crore</td></tr>
    <tr><td>Interest Rate</td><td>{cfg.get('interest_rate', 0.115)*100:.1f}% p.a.</td></tr>
    <tr><td>EMI</td><td>Rs {cfg.get('emi_lac_mth', 0):.2f} Lac/month</td></tr>
    <tr><td>Selling Price</td><td>Rs {cfg.get('selling_price_per_mt', 35000):,}/MT</td></tr>
    <tr><td>Profit/MT</td><td>Rs {cfg.get('profit_per_mt', 0):,}</td></tr>
    <tr><td>Monthly Profit</td><td>Rs {cfg.get('monthly_profit_lac', 0):.1f} Lac</td></tr>
    </table>

    <h2>6. Project Timeline</h2>
    <p>Total duration: 12-18 months from DPR approval to commercial production.</p>

    <p style="text-align:center; margin-top:40px; color:#666;">
    CONFIDENTIAL — For Private Circulation Only<br>
    {company.get('name', 'PPS Anantams Corporation Pvt Ltd')} | GST: {company.get('gst', '')}
    </p>
    </body></html>
    """

    # Convert to PDF-like output (save as HTML for now — PDF needs wkhtmltopdf)
    buf = io.BytesIO(html.encode("utf-8"))
    filename = f"DPR_{int(tpd)}TPD_{state}_{datetime.now().strftime('%Y%m%d')}.html"

    try:
        from engines.audit_logger import register_export, log_event
        _ensure()
        fpath = EXPORT_DIR / filename
        fpath.write_bytes(buf.getvalue())
        buf.seek(0)
        register_export("dpr", str(fpath), cfg)
        log_event("EXPORT_CREATED", "report_generator", f"DPR {filename}", "success")
    except Exception:
        pass

    return buf, filename
