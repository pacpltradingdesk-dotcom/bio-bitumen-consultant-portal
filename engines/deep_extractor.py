"""
Bio Bitumen — DEEP Content Extractor
======================================
Reads EVERY file type: DOCX, PDF, PPTX, XLSX, PNG
Extracts ALL content: text, tables, numbers, slide content
Builds a unified data store that powers the entire dashboard.
"""
import os
import sys
import json
import time
import traceback
from pathlib import Path

DOC_ROOT = Path(r"C:\Users\HP\Desktop\Bio Bitumen Full Working all document")
DATA_DIR = Path(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))) / "data"
DEEP_CACHE = DATA_DIR / "deep_content.json"
SKIP_DIRS = {'__pycache__', '.claude', 'consultant_portal', 'CUSTOMER_PACKAGES',
             'ARCHIVE_Output_100Cr_Original', 'ARCHIVE_Output_Biomass_Original',
             'ARCHIVE_Output_Bio_Bitumen_Original'}


def extract_docx(filepath):
    """Extract ALL text + tables from Word document."""
    try:
        from docx import Document
        doc = Document(filepath)
        text_parts = []
        tables_data = []

        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text.strip())

        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    rows.append(cells)
            if rows:
                tables_data.append(rows)

        return {
            "text": "\n".join(text_parts),
            "tables": tables_data,
            "paragraphs": len(text_parts),
            "table_count": len(tables_data),
        }
    except Exception as e:
        return {"text": "", "error": str(e)}


def extract_pdf(filepath):
    """Extract text + tables from PDF."""
    try:
        import pdfplumber
        text_parts = []
        tables_data = []

        with pdfplumber.open(filepath) as pdf:
            for i, page in enumerate(pdf.pages[:20]):  # Max 20 pages
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
                page_tables = page.extract_tables()
                if page_tables:
                    for t in page_tables:
                        clean = [[str(c or '') for c in row] for row in t if any(c for c in row)]
                        if clean:
                            tables_data.append(clean)

        return {
            "text": "\n".join(text_parts),
            "tables": tables_data,
            "pages": min(len(text_parts), 20),
            "table_count": len(tables_data),
        }
    except Exception as e:
        return {"text": "", "error": str(e)}


def extract_pptx(filepath):
    """Extract ALL slide content from PowerPoint."""
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        slides_data = []

        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_text.append(para.text.strip())
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        if any(cells):
                            slide_text.append(" | ".join(cells))
            if slide_text:
                slides_data.append({
                    "slide": i + 1,
                    "content": "\n".join(slide_text),
                })

        return {
            "text": "\n\n".join(s["content"] for s in slides_data),
            "slides": slides_data,
            "slide_count": len(slides_data),
        }
    except Exception as e:
        return {"text": "", "error": str(e)}


def extract_xlsx(filepath):
    """Extract ALL sheet data from Excel."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(filepath, data_only=True)
        sheets_data = {}

        for name in wb.sheetnames:
            ws = wb[name]
            rows = []
            for row in ws.iter_rows(min_row=1, max_row=min(50, ws.max_row), values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    rows.append(cells)
            sheets_data[name] = rows

        wb.close()
        text = "\n".join(
            f"[{name}] " + " | ".join(row)
            for name, rows in sheets_data.items()
            for row in rows[:10]
        )
        return {
            "text": text,
            "sheets": sheets_data,
            "sheet_count": len(sheets_data),
        }
    except Exception as e:
        return {"text": "", "error": str(e)}


def extract_image_info(filepath):
    """Extract image dimensions and basic info."""
    try:
        from PIL import Image
        img = Image.open(filepath)
        return {
            "text": f"Image: {img.size[0]}x{img.size[1]} pixels, {img.mode}",
            "width": img.size[0],
            "height": img.size[1],
            "format": img.format,
        }
    except Exception as e:
        return {"text": "", "error": str(e)}


# Extension → extractor mapping
EXTRACTORS = {
    ".docx": extract_docx,
    ".pdf": extract_pdf,
    ".pptx": extract_pptx,
    ".xlsx": extract_xlsx,
    ".png": extract_image_info,
    ".jpg": extract_image_info,
    ".jpeg": extract_image_info,
}


def deep_scan_all(max_files=None, progress_callback=None):
    """
    Deep scan ALL files in the project.
    Returns a dict: {relative_path: {filename, ext, content, tables, ...}}
    """
    results = {}
    file_count = 0
    error_count = 0
    skipped = 0

    for root, dirs, files in os.walk(str(DOC_ROOT)):
        dirs[:] = [d for d in dirs if d not in SKIP_DIRS]
        rel_dir = os.path.relpath(root, str(DOC_ROOT))

        for fname in sorted(files):
            ext = os.path.splitext(fname)[1].lower()
            if ext not in EXTRACTORS:
                skipped += 1
                continue

            filepath = os.path.join(root, fname)
            rel_path = os.path.relpath(filepath, str(DOC_ROOT)).replace("\\", "/")

            # Skip very large files (>5MB for DOCX/PDF, >10MB for XLSX)
            try:
                size = os.path.getsize(filepath)
                if ext in ('.docx', '.pdf') and size > 5_000_000:
                    skipped += 1
                    continue
                if ext == '.xlsx' and size > 10_000_000:
                    skipped += 1
                    continue
            except:
                continue

            file_count += 1
            if max_files and file_count > max_files:
                break

            try:
                extractor = EXTRACTORS[ext]
                content = extractor(filepath)

                results[rel_path] = {
                    "filename": fname,
                    "extension": ext,
                    "path": rel_path,
                    "size": size,
                    "content": content.get("text", "")[:2000],  # Cap at 2000 chars
                    "tables_count": content.get("table_count", content.get("sheet_count", 0)),
                    "has_tables": bool(content.get("tables") or content.get("sheets")),
                    "slide_count": content.get("slide_count", 0),
                    "pages": content.get("pages", 0),
                    "error": content.get("error", ""),
                    # Capacity detection
                    "capacity": _detect_capacity(rel_path),
                    "section": _detect_section(rel_path),
                }

                if progress_callback and file_count % 50 == 0:
                    progress_callback(file_count, fname)

            except Exception as e:
                error_count += 1
                results[rel_path] = {
                    "filename": fname, "extension": ext, "path": rel_path,
                    "content": "", "error": str(e),
                }

        if max_files and file_count > max_files:
            break

    return results, file_count, error_count, skipped


def _detect_capacity(path):
    for cap in ["05MT", "10MT", "15MT", "20MT", "30MT", "40MT", "50MT"]:
        if cap in path:
            return cap
    return "general"


def _detect_section(path):
    parts = path.replace("\\", "/").split("/")
    for part in parts:
        if part[:2].isdigit() and "_" in part:
            return part
    return ""


def save_deep_scan(results):
    """Save results to cache file."""
    os.makedirs(str(DATA_DIR), exist_ok=True)
    # Save without full content to keep file manageable
    slim = {}
    for path, data in results.items():
        slim[path] = {
            "filename": data["filename"],
            "extension": data["extension"],
            "capacity": data.get("capacity", ""),
            "section": data.get("section", ""),
            "content_length": len(data.get("content", "")),
            "has_tables": data.get("has_tables", False),
            "tables_count": data.get("tables_count", 0),
            "slide_count": data.get("slide_count", 0),
            "pages": data.get("pages", 0),
            "error": data.get("error", ""),
            "summary": data.get("content", "")[:300],
        }
    with open(str(DEEP_CACHE), "w", encoding="utf-8") as f:
        json.dump(slim, f, indent=2, ensure_ascii=False)
    return str(DEEP_CACHE)


def load_deep_scan():
    """Load cached deep scan results."""
    if DEEP_CACHE.exists():
        with open(str(DEEP_CACHE), "r", encoding="utf-8") as f:
            return json.load(f)
    return {}


def get_scan_stats():
    """Get statistics from the last deep scan."""
    data = load_deep_scan()
    if not data:
        return None

    stats = {
        "total_files": len(data),
        "by_type": {},
        "by_capacity": {},
        "with_content": 0,
        "with_tables": 0,
        "with_errors": 0,
    }
    for path, entry in data.items():
        ext = entry.get("extension", "")
        cap = entry.get("capacity", "general")
        stats["by_type"][ext] = stats["by_type"].get(ext, 0) + 1
        stats["by_capacity"][cap] = stats["by_capacity"].get(cap, 0) + 1
        if entry.get("content_length", 0) > 0:
            stats["with_content"] += 1
        if entry.get("has_tables"):
            stats["with_tables"] += 1
        if entry.get("error"):
            stats["with_errors"] += 1

    return stats
