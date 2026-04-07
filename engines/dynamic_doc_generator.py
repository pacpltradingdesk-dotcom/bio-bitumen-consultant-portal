"""
Bio Bitumen — DYNAMIC DOCUMENT GENERATOR
==========================================
Generates fresh DOCX, PPTX, XLSX for ANY capacity (5-100 TPD).
Uses current dashboard config — change anything, documents auto-update.
Covers: DPR, Financial Model, Bank Proposal, Investor Pitch, Compliance,
        HR Plan, BOQ, Technical Report, Government Forms.
"""
import os
import io
from datetime import datetime, timezone, timedelta
from docx import Document
from docx.shared import Inches, Pt, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import nsdecls
from docx.oxml import parse_xml

IST = timezone(timedelta(hours=5, minutes=30))
NAVY = RGBColor(0x00, 0x33, 0x66)
TEAL = RGBColor(0x00, 0x66, 0x99)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _add_heading(doc, text, level=1):
    h = doc.add_heading(text, level=level)
    for run in h.runs:
        run.font.color.rgb = NAVY if level == 1 else TEAL
    return h


def _add_table(doc, headers, rows):
    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
    table.style = 'Light Grid Accent 1'
    # Headers
    for i, h in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = str(h)
        for p in cell.paragraphs:
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(10)
    # Data rows
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            table.rows[r_idx + 1].cells[c_idx].text = str(val)
    return table


def _footer(doc, company):
    for section in doc.sections:
        footer = section.footer
        footer.is_linked_to_previous = False
        p = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
        p.text = f"{company['name']} | {company['owner']} | {company['phone']} | CONFIDENTIAL"
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in p.runs:
            run.font.size = Pt(8)
            run.font.color.rgb = RGBColor(0x88, 0x88, 0x88)


# ═══════════════════════════════════════════════════════════════════════
# 1. DPR (DETAILED PROJECT REPORT) — DOCX
# ═══════════════════════════════════════════════════════════════════════
def generate_dpr_docx(cfg, company, customer=None):
    """Generate complete DPR as DOCX using current dashboard config + client info."""
    doc = Document()
    now = datetime.now(IST)
    p = cfg.get("plant_data", {})

    # Client info — from cfg (Project Setup) or customer parameter
    client_name = cfg.get("client_name") or (customer.get("name", "") if customer else "") or "The Promoter"
    client_company = cfg.get("client_company") or (customer.get("company", "") if customer else "") or ""
    project_name = cfg.get("project_name") or f"Bio-Modified Bitumen Plant — {cfg.get('capacity_tpd', 20):.0f} MT/Day"
    site_address = cfg.get("site_address") or f"{cfg.get('location', 'To be finalized')}, {cfg.get('state', '')}"

    # Cover Page
    doc.add_paragraph("")
    title = doc.add_heading("DETAILED PROJECT REPORT", level=0)
    for run in title.runs:
        run.font.color.rgb = NAVY

    doc.add_heading(project_name, level=1)

    if client_company:
        doc.add_paragraph(f"Prepared for: {client_name}")
        doc.add_paragraph(f"Company: {client_company}")
    elif client_name and client_name != "The Promoter":
        doc.add_paragraph(f"Prepared for: {client_name}")

    doc.add_paragraph(f"Total Investment: Rs {cfg.get('investment_cr', 8):.2f} Crore")
    doc.add_paragraph(f"Plant Capacity: {cfg.get('capacity_tpd', 20):.0f} MT/Day")
    doc.add_paragraph(f"Location: {site_address}")
    if cfg.get("site_pincode"):
        doc.add_paragraph(f"Pincode: {cfg.get('site_pincode', "")} | District: {cfg.get('site_district', '')}")
    if cfg.get("site_area_acres"):
        doc.add_paragraph(f"Site Area: {cfg.get('site_area_acres', 0)} Acres ({cfg.get('site_ownership', '')})")
    doc.add_paragraph("")
    doc.add_paragraph(f"Prepared by: {company['trade_name']}")
    doc.add_paragraph(f"Consultant: {company['owner']} | {company['phone']}")
    doc.add_paragraph(f"Date: {now.strftime('%d %B %Y')}")
    if cfg.get("project_id"):
        doc.add_paragraph(f"Reference: {cfg.get('project_id', "")}")
    doc.add_paragraph(f"CONFIDENTIAL — For {client_name} use only")
    doc.add_page_break()

    # 1. Executive Summary
    _add_heading(doc, "1. EXECUTIVE SUMMARY")
    doc.add_paragraph(
        f"This Detailed Project Report presents a {cfg.get('capacity_tpd', 20):.0f} MT/Day Bio-Modified Bitumen "
        f"manufacturing plant using agricultural biomass pyrolysis technology. The project requires a total "
        f"investment of Rs {cfg.get('investment_cr', 8):.2f} Crore with a debt:equity ratio of "
        f"{int((1-cfg.get('equity_ratio', 0.40))*100)}:{int(cfg.get('equity_ratio', 0.40)*100)}."
    )
    _add_table(doc, ["Parameter", "Value"], [
        ["Plant Capacity", f"{cfg.get('capacity_tpd', 20):.0f} MT/Day"],
        ["Total Investment", f"Rs {cfg.get('investment_cr', 8):.2f} Crore"],
        ["Bank Loan", f"Rs {cfg.get('loan_cr', 4.8):.2f} Crore"],
        ["Promoter Equity", f"Rs {cfg.get('equity_cr', 3.2):.2f} Crore"],
        ["Annual Revenue (Yr5)", f"Rs {cfg.get('revenue_yr5_lac', 0):.0f} Lakhs"],
        ["ROI", f"{cfg.get('roi_pct', 0):.1f}%"],
        ["IRR", f"{cfg.get('irr_pct', 0):.1f}%"],
        ["DSCR (Year 3)", f"{cfg.get('dscr_yr3', 0):.2f}x"],
        ["Break-Even", f"{cfg.get('break_even_months', 0)} months"],
        ["Monthly EMI", f"Rs {cfg.get('emi_lac_mth', 0):.2f} Lakhs"],
        ["Staff Requirement", f"{cfg.get('staff', 20)} persons"],
        ["Power Requirement", f"{cfg.get('power_kw', 100):.0f} kW"],
    ])
    doc.add_page_break()

    # 2. Market Analysis
    _add_heading(doc, "2. MARKET ANALYSIS")
    doc.add_paragraph(
        "India is the third-largest road network in the world with over 6.4 million km. "
        "The Government of India has allocated Rs 2.7 lakh crore for highway development in FY 2025-26. "
        "Bitumen demand is projected to grow at 8-10% CAGR driven by Bharatmala, PMGSY, and state highway projects."
    )
    doc.add_paragraph(
        f"Bio-Modified Bitumen offers a sustainable alternative at Rs {cfg.get('selling_price_per_mt', 35000):,}/MT "
        f"compared to conventional VG30 at Rs 48,000-52,000/MT, providing a significant cost advantage "
        f"while meeting NHAI green mandate requirements."
    )

    # 2A. PROJECT JUSTIFICATION & RATIONALE
    doc.add_heading("2.1 Project Justification & Rationale", level=2)
    doc.add_paragraph(
        f"India imports 49% of its bitumen requirement (Rs 25,000+ Crore/year). "
        f"The Government has mandated 15% bio-bitumen usage by 2030 through NHAI. "
        f"130-216 new bio-bitumen plants are needed in the next 5-7 years. "
        f"This project addresses both import substitution and sustainable road construction."
    )

    # 2B. DEMAND-SUPPLY GAP
    doc.add_heading("2.2 Demand-Supply Gap Analysis", level=2)
    ann_prod = cfg.get('capacity_tpd', 20) * cfg.get('working_days', 300) * 0.4
    _add_table(doc, ["Parameter", "Value"], [
        ["India Annual Bitumen Demand", "8.8 Million MT"],
        ["India Annual Bitumen Production", "4.5 Million MT"],
        ["Import Dependency", "49% (~4.3 Million MT)"],
        ["Bio-Bitumen Target (NHAI 15%)", "1.32 Million MT by 2030"],
        ["Current Bio-Bitumen Supply", "< 10,000 MT (negligible)"],
        [f"This Plant Annual Output", f"{ann_prod:,.0f} MT/year"],
    ])

    # 2C. TARGET CUSTOMERS
    doc.add_heading("2.3 Target Customer Segments", level=2)
    doc.add_paragraph(
        "Primary Buyers: NHAI contractors, State PWD contractors, Municipal corporations. "
        "Secondary: Private road developers, airport authorities, port trusts. "
        "By-Products: Farmers (biochar), industrial boilers (bio-oil), carbon market (credits)."
    )

    # 2D. COMPETITION ANALYSIS
    doc.add_heading("2.4 Competition Analysis", level=2)
    _add_table(doc, ["Competitor", "Capacity", "Technology", "Our Advantage"], [
        ["Praj Industries", "Pilot scale", "Lignocellulosic", "We are production-ready, they are pilot"],
        ["Indian Oil (IOCL)", "In-house R&D", "Refinery-based", "We use agro-waste (lower cost)"],
        ["Local pyrolysis units", "5-10 TPD", "Basic", "We have CSIR-CRRI license + IS:73 compliance"],
        ["Import (Iran/Singapore)", "Bulk", "Conventional", "We eliminate import dependency"],
    ])

    # 2E. MARKETING STRATEGY
    doc.add_heading("2.5 Marketing Strategy & Sales Plan", level=2)
    doc.add_paragraph(
        "1. GeM Portal registration for direct government supply\n"
        "2. NHAI contractor partnerships (tie-up with 5-10 major contractors)\n"
        "3. State PWD empanelment for each target state\n"
        "4. Direct sales to private road developers\n"
        "5. Bio-char sales through agricultural cooperatives\n"
        "6. Carbon credit registration with voluntary market (Verra/Gold Standard)"
    )

    # 2F. REVENUE MODEL
    doc.add_heading("2.6 Revenue Model & Pricing Strategy", level=2)
    _add_table(doc, ["Product", "Price (Rs/MT)", "Share of Revenue"], [
        ["Bio-Bitumen VG30", f"{cfg.get('sale_bio_bitumen_vg30', 44000):,}", "60-70%"],
        ["Bio-Bitumen VG40", f"{cfg.get('sale_bio_bitumen_vg40', 48000):,}", "15-20%"],
        ["Bio-Char (Agriculture)", f"{cfg.get('sale_biochar_agri', 26000):,}", "8-12%"],
        ["Bio-Oil (Fuel)", f"{cfg.get('sale_bio_oil_fuel', 22000):,}", "3-5%"],
        ["Carbon Credits", f"{cfg.get('sale_carbon_credit', 12500):,}/credit", "2-3%"],
    ])
    doc.add_page_break()

    # 3. Technical Details
    _add_heading(doc, "3. TECHNICAL DETAILS")
    doc.add_heading("3.1 Process Flow", level=2)
    doc.add_paragraph(
        "Biomass Collection → Size Reduction (Shredder) → Drying (<15% moisture) → "
        "Pyrolysis (400-600°C, oxygen-free) → Bio-Oil (40%) + Biochar (30%) + Syngas (25%) → "
        "Bio-Oil blended with bitumen → Bio-Modified Bitumen → Quality Testing → Dispatch"
    )
    doc.add_heading("3.2 Plant Specifications", level=2)
    _add_table(doc, ["Specification", "Value"], [
        ["Capacity", f"{cfg.get('capacity_tpd', 20):.0f} MT/Day biomass input"],
        ["Bio-Oil Output", f"{cfg.get('oil_ltr_day', 8000):.0f} Litres/day"],
        ["Biochar Output", f"{cfg.get('char_kg_day', 6000):.0f} Kg/day"],
        ["Power Requirement", f"{cfg.get('power_kw', 100):.0f} kW"],
        ["Water Requirement", f"{max(5, int(cfg.get('capacity_tpd', 20))):,} KLD"],
        ["Land Requirement", f"{int(cfg.get('capacity_tpd', 20) * 520):,} sq ft ({cfg.get('capacity_tpd', 20) * 520 / 43560:.2f} acres)"],
        ["Manpower", f"{cfg.get('staff', 20)} persons"],
        ["Working Days", f"{cfg.get('working_days', 300)}/year"],
    ])
    doc.add_page_break()

    # 4. Investment Breakdown
    _add_heading(doc, "4. INVESTMENT BREAKDOWN")
    inv_rows = [
        ["Civil & Building", f"Rs {p.get('civil_lac', 0):.1f} Lac"],
        ["Machinery & Equipment", f"Rs {p.get('mach_lac', 0):.1f} Lac"],
        ["GST on Machinery (18%)", f"Rs {p.get('gst_mach_lac', 0):.1f} Lac"],
        ["Working Capital", f"Rs {p.get('wc_lac', 0):.1f} Lac"],
        ["Interest During Construction", f"Rs {p.get('idc_lac', 0):.1f} Lac"],
        ["Pre-operative Expenses", f"Rs {p.get('preop_lac', 0):.1f} Lac"],
        ["Contingency", f"Rs {p.get('cont_lac', 0):.1f} Lac"],
        ["Security Deposit", f"Rs {p.get('sec_lac', 0):.1f} Lac"],
        ["TOTAL", f"Rs {cfg.get('investment_lac', 800):.1f} Lac (Rs {cfg.get('investment_cr', 8):.2f} Crore)"],
    ]
    _add_table(doc, ["Component", "Amount"], inv_rows)

    # 5. Financial Projections
    _add_heading(doc, "5. FINANCIAL PROJECTIONS (7 YEARS)")
    if cfg.get("roi_timeline"):
        headers = ["Year", "Revenue (Lac)", "EBITDA (Lac)", "PAT (Lac)", "DSCR"]
        rows = []
        for yr in cfg.get("roi_timeline", 0):
            rows.append([
                str(yr["Year"]),
                f"{yr['Revenue (Lac)']:.0f}",
                f"{yr['EBITDA (Lac)']:.0f}",
                f"{yr['PAT (Lac)']:.0f}",
                f"{yr['DSCR']:.2f}x",
            ])
        _add_table(doc, headers, rows)

    doc.add_paragraph("")
    doc.add_paragraph(f"Profit per MT: Rs {cfg.get('profit_per_mt', 0):,.0f}")
    doc.add_paragraph(f"Monthly Profit (Yr5): Rs {cfg.get('monthly_profit_lac', 0):.1f} Lakhs")

    # 6. Cost Structure
    _add_heading(doc, "6. COST STRUCTURE (Per MT Output)")
    _add_table(doc, ["Cost Item", "Rs/MT"], [
        ["Raw Material (Biomass)", f"{cfg.get('raw_material_cost_per_mt', 8000):,}"],
        ["Power & Fuel", f"{cfg.get('power_cost_per_mt', 4500):,}"],
        ["Labour", f"{cfg.get('labour_cost_per_mt', 3000):,}"],
        ["Chemicals & Consumables", f"{cfg.get('chemical_cost_per_mt', 1500):,}"],
        ["Packaging & Handling", f"{cfg.get('packaging_cost_per_mt', 500):,}"],
        ["Transport (Outbound)", f"{cfg.get('transport_cost_per_mt', 2000):,}"],
        ["QC & Testing", f"{cfg.get('qc_cost_per_mt', 500):,}"],
        ["Miscellaneous", f"{cfg.get('misc_cost_per_mt', 1000):,}"],
        ["TOTAL VARIABLE COST", f"{cfg.get('total_variable_cost_per_mt', 22000):,}"],
        ["Selling Price", f"{cfg.get('selling_price_per_mt', 35000):,}"],
        ["PROFIT PER MT", f"{cfg.get('profit_per_mt', 0):,}"],
    ])

    # 7. Risk Analysis
    _add_heading(doc, "7. RISK ANALYSIS")
    _add_table(doc, ["Risk", "Mitigation"], [
        ["Raw material price increase", "Long-term FPO contracts, multi-crop sourcing"],
        ["Bitumen price volatility", "Forward contracts, diversified product mix"],
        ["Technology risk", "Proven pyrolysis tech, CSIR-CRRI validated"],
        ["Demand risk", "Govt mandate for green materials, NHAI preference"],
        ["Regulatory risk", "All licenses identified, compliance calendar maintained"],
        ["Monsoon seasonality", "3-month raw material buffer storage"],
    ])

    # 8. Compliance
    _add_heading(doc, "8. LICENSES & COMPLIANCE")
    doc.add_paragraph("All required licenses identified with typical timelines:")
    _add_table(doc, ["License", "Authority", "Timeline"], [
        ["GST Registration", "Central GST", "7 days"],
        ["Factory License", "State Inspector", "30 days"],
        ["Consent to Establish (CTE)", "State PCB", "60 days"],
        ["Fire NOC", "Fire Department", "30 days"],
        ["PESO License", "PESO Nagpur", "90 days"],
        ["Udyam (MSME)", "MSME Ministry", "1 day"],
        ["EPF/ESI Registration", "EPFO/ESIC", "7 days"],
        ["Electricity HT Connection", "State DISCOM", "60 days"],
    ])

    doc.add_page_break()

    # 9. BOQ — Bill of Quantities (auto-generated from capacity)
    _add_heading(doc, "9. BILL OF QUANTITIES (BOQ)")
    doc.add_paragraph(f"Equipment and material requirements for {cfg.get('capacity_tpd', 20):.0f} MT/Day plant:")
    try:
        from state_manager import calculate_boq
        boq = calculate_boq(cfg.get("capacity_tpd", 0))
        boq_rows = []
        total_boq = 0
        for item in boq:
            boq_rows.append([item["item"], item["spec"], str(item["qty"]) + " " + item["unit"],
                              f"Rs {item['rate_lac']:.1f} Lac", f"Rs {item['amount_lac']:.1f} Lac"])
            total_boq += item["amount_lac"]
        boq_rows.append(["TOTAL", "", "", "", f"Rs {total_boq:.1f} Lac (Rs {total_boq/100:.2f} Cr)"])
        _add_table(doc, ["Equipment", "Specification", "Qty", "Rate", "Amount"], boq_rows)
    except Exception:
        doc.add_paragraph("BOQ will be provided separately based on final design.")
    doc.add_page_break()

    # 10. Raw Material Sourcing
    _add_heading(doc, "10. RAW MATERIAL SOURCING PLAN")
    doc.add_paragraph(
        f"The plant requires approximately {cfg.get('capacity_tpd', 20)*2.5:.0f} MT/day of agricultural biomass. "
        f"Primary source: {cfg.get('biomass_source', 'Rice Straw')} from local farmer cooperatives "
        f"within 50-100 km radius of {cfg.get('location', 'the plant site')}."
    )
    doc.add_heading("10.1 Biomass Types & Yield", level=2)
    _add_table(doc, ["Biomass", "Oil Yield %", "Char Yield %", "Cost Rs/MT"], [
        ["Rice Straw", "40%", "32%", "1,200-1,500"],
        ["Wheat Straw", "38%", "30%", "1,400-1,600"],
        ["Sugarcane Bagasse", "42%", "28%", "1,000-1,200"],
        ["Cotton Stalk", "35%", "33%", "1,600-2,000"],
        ["Groundnut Shell", "38%", "35%", "2,200-2,500"],
    ])
    doc.add_heading("10.2 Seasonal Procurement Calendar", level=2)
    _add_table(doc, ["Quarter", "Primary Crop", "Availability", "Strategy"], [
        ["Q1 (Apr-Jun)", "Wheat Straw, Sugarcane Bagasse", "High", "Peak procurement + pelletization"],
        ["Q2 (Jul-Sep)", "Monsoon Season", "Low", "Use pelletized stock from buffer storage"],
        ["Q3 (Oct-Dec)", "Rice Straw, Cotton Stalk", "Very High", "Maximum procurement, 90-day buffer"],
        ["Q4 (Jan-Mar)", "Mixed sources", "Medium", "Supplementary procurement"],
    ])
    doc.add_paragraph(f"90-Day Buffer Stock Required: {cfg.get('capacity_tpd', 20)*2.5*90:.0f} MT")
    doc.add_paragraph(f"Storage Area Required: {int(cfg.get('capacity_tpd', 20)*2.5*90/80*3.28*3.28/3):,} sq ft")
    doc.add_page_break()

    # 11. Environmental Impact
    _add_heading(doc, "11. ENVIRONMENTAL IMPACT & CARBON CREDITS")
    co2_annual = cfg.get("capacity_tpd", 0) * 300 * 0.35
    stubble = cfg.get("capacity_tpd", 0) * 300 * 2.5
    carbon_lac = co2_annual * 12 * 84 / 100000
    doc.add_paragraph(
        f"The Bio-Modified Bitumen plant provides significant environmental benefits by diverting "
        f"agricultural waste from open burning and reducing petroleum dependency."
    )
    _add_table(doc, ["Environmental Metric", "Annual Impact"], [
        ["CO2 Emissions Avoided", f"{co2_annual:,.0f} tonnes/year"],
        ["Crop Residue Diverted", f"{stubble:,.0f} MT/year"],
        ["Carbon Credit Revenue", f"Rs {carbon_lac:.1f} Lakhs/year (at USD 12/tonne)"],
        ["Equivalent Trees Planted", f"{co2_annual*50:,.0f} trees equivalent"],
        ["Equivalent Cars Removed", f"{co2_annual/4.6:,.0f} cars/year"],
        ["Water Saved vs Conventional", "15% less usage"],
        ["NHAI Green Mandate Compliance", "YES — 15% bio-bitumen target by 2030"],
    ])
    doc.add_page_break()

    # 12. Government Subsidies
    _add_heading(doc, "12. GOVERNMENT SUBSIDIES & INCENTIVES")
    state = cfg.get("state", "Gujarat")
    subsidy_pct = 20  # Default
    doc.add_paragraph(
        f"The project is eligible for multiple government subsidies and incentives. "
        f"Located in {state}, the following schemes are applicable:"
    )
    _add_table(doc, ["Scheme", "Authority", "Benefit", "Status"], [
        ["MNRE Waste-to-Wealth", "Ministry of New & Renewable Energy", "25% capital subsidy", "Eligible"],
        [f"{state} MSME Subsidy", f"{state} Industrial Policy", f"{subsidy_pct}% capital subsidy", "Eligible"],
        ["CGTMSE Guarantee", "CGTMSE Trust", "Collateral-free loan up to Rs 5 Cr", "Eligible"],
        ["Carbon Credits", "Voluntary Carbon Market", f"Rs {carbon_lac:.1f} Lac/year", "Applicable"],
        ["Udyam Registration", "MSME Ministry", "Priority sector lending", "Eligible"],
        ["Technology Development Board", "DST", "Up to 50% for new technology", "Applicable"],
    ])
    subsidy_amount = cfg.get("investment_cr", 0) * subsidy_pct / 100
    doc.add_paragraph(
        f"\nTotal Estimated Subsidy: Rs {subsidy_amount:.2f} Crore ({subsidy_pct}% of investment)"
    )
    doc.add_paragraph(
        f"Effective Investment After Subsidy: Rs {cfg.get('investment_cr', 8) - subsidy_amount:.2f} Crore"
    )
    doc.add_page_break()

    # 13. Implementation Timeline
    _add_heading(doc, "13. IMPLEMENTATION TIMELINE (12-18 MONTHS)")
    _add_table(doc, ["Phase", "Duration", "Activities"], [
        ["1. Pre-Feasibility & DPR", "1 month", "Site selection, feasibility study, DPR preparation"],
        ["2. Company Setup", "1 month", "ROC incorporation, PAN, TAN, GST, Udyam registration"],
        ["3. Land & Approvals", "3 months", "Land purchase/lease, building plan, zoning approval"],
        ["4. Environmental Clearances", "5 months", "CTE from PCB, EIA, fire NOC, PESO license"],
        ["5. Bank Loan Sanction", "3 months", "DPR submission, CMA data, bank appraisal"],
        ["6. Engineering & Design", "3 months", "Detailed engineering, P&ID, civil design, electrical SLD"],
        ["7. Equipment Procurement", "4 months", "Equipment ordering, vendor finalization, delivery tracking"],
        ["8. Civil Construction", "5 months", "Foundation, structure, building, roads, bund walls"],
        ["9. Installation & Commissioning", "3 months", "Equipment erection, piping, electrical, trial run"],
        ["10. Commercial Production", "6 months", "CTO from PCB, ramp-up: 40% to 85% utilization"],
    ])
    doc.add_page_break()

    # 14. Manpower Plan
    _add_heading(doc, "14. MANPOWER PLAN")
    _add_table(doc, ["Role", "Count", "Monthly CTC (Rs)"], [
        ["Plant Manager", "1", "60,000-80,000"],
        ["Shift Supervisor", "2", "30,000-40,000"],
        ["Reactor Operators", f"{max(4, int(cfg.get('capacity_tpd', 20)*0.3))}", "18,000-25,000"],
        ["Lab Technician", "1-2", "20,000-30,000"],
        ["Electrician", "1", "20,000-25,000"],
        ["Helpers/Loaders", f"{max(6, int(cfg.get('capacity_tpd', 20)*0.5))}", "12,000-15,000"],
        ["Drivers", "2-3", "15,000-20,000"],
        ["Office/Accounts", "2", "20,000-30,000"],
        ["Security", "2-3", "12,000-15,000"],
        ["TOTAL", f"{cfg.get('staff', 20)}", f"Rs {cfg.get('payroll_lac_yr', 45):.1f} Lac/year"],
    ])
    doc.add_page_break()

    # 15. Sensitivity Analysis
    _add_heading(doc, "15. SENSITIVITY ANALYSIS")
    doc.add_paragraph("The following table shows project profitability under different scenarios:")
    if cfg.get("sensitivity_matrix"):
        _add_table(doc, ["Scenario", "Low Price", "Base Price", "High Price"], [
            ["Low Cost", f"Rs {cfg.get('sensitivity_matrix', [])[0][0]:.0f} Lac",
             f"Rs {cfg.get('sensitivity_matrix', [])[0][1]:.0f} Lac", f"Rs {cfg.get('sensitivity_matrix', [])[0][2]:.0f} Lac"],
            ["Base Cost", f"Rs {cfg.get('sensitivity_matrix', [])[1][0]:.0f} Lac",
             f"Rs {cfg.get('sensitivity_matrix', [])[1][1]:.0f} Lac", f"Rs {cfg.get('sensitivity_matrix', [])[1][2]:.0f} Lac"],
            ["High Cost", f"Rs {cfg.get('sensitivity_matrix', [])[2][0]:.0f} Lac",
             f"Rs {cfg.get('sensitivity_matrix', [])[2][1]:.0f} Lac", f"Rs {cfg.get('sensitivity_matrix', [])[2][2]:.0f} Lac"],
        ])
    doc.add_paragraph(
        f"Even in the worst-case scenario (high cost + low price), the project maintains positive EBITDA "
        f"due to the inherent cost advantage of bio-modified bitumen vs conventional petroleum bitumen."
    )
    doc.add_page_break()

    # 16. Consultant Profile
    # 15A. NPV CALCULATION
    doc.add_heading("15.1 Net Present Value (NPV)", level=2)
    npv = cfg.get('npv_lac', 0)
    doc.add_paragraph(
        f"NPV at 15% discount rate: Rs {npv:.1f} Lac. "
        f"A positive NPV indicates the project generates returns above the cost of capital."
    )

    # 15B. PROJECTED BALANCE SHEET
    doc.add_heading("15.2 Projected Balance Sheet Summary", level=2)
    if cfg.get('balance_sheet'):
        bs_rows = [[item.get('Item', ''), f"Rs {item.get('Amount (Lac)', 0):.1f} Lac"]
                    for item in cfg['balance_sheet']]
        _add_table(doc, ["Item", "Amount"], bs_rows)
    else:
        inv_lac = cfg.get('investment_cr', 8) * 100
        doc.add_paragraph(
            f"Total Assets (Year 5): Rs {inv_lac * 0.7:.0f} Lac (net of depreciation)\n"
            f"Net Worth (Year 5): Rs {cfg.get('net_worth_yr5_lac', inv_lac * 0.5):.0f} Lac"
        )

    # 15C. TERM LOAN REPAYMENT SCHEDULE
    doc.add_heading("15.3 Term Loan Repayment Schedule", level=2)
    loan_lac = cfg.get('loan_cr', 6) * 100
    emi = cfg.get('emi_lac_mth', 7)
    tenure = cfg.get('emi_tenure_months', 84)
    doc.add_paragraph(
        f"Loan Amount: Rs {loan_lac:.0f} Lac | EMI: Rs {emi:.2f} Lac/month | "
        f"Tenure: {tenure} months ({tenure//12} years) | "
        f"Interest Rate: {cfg.get('interest_rate', 0.115)*100:.1f}% p.a."
    )
    _add_table(doc, ["Year", "Opening Balance", "Interest", "Principal", "Closing Balance"], [
        [f"Year {y}", f"Rs {max(0, loan_lac - loan_lac/7*(y-1)):.0f}L",
         f"Rs {max(0, loan_lac - loan_lac/7*(y-1))*cfg.get('interest_rate', 0.115):.0f}L",
         f"Rs {loan_lac/7:.0f}L",
         f"Rs {max(0, loan_lac - loan_lac/7*y):.0f}L"]
        for y in range(1, 8)
    ])

    # 15D. COST OF PRODUCTION PER UNIT
    doc.add_heading("15.4 Cost of Production per MT", level=2)
    _add_table(doc, ["Cost Head", "Rs/MT"], [
        ["Raw Material", f"{cfg.get('raw_material_cost_per_mt', 8000):,}"],
        ["Power & Fuel", f"{cfg.get('power_cost_per_mt', 4500):,}"],
        ["Labour", f"{cfg.get('labour_cost_per_mt', 3000):,}"],
        ["Chemicals", f"{cfg.get('chemical_cost_per_mt', 1500):,}"],
        ["Packing", f"{cfg.get('packaging_cost_per_mt', 500):,}"],
        ["Transport", f"{cfg.get('transport_cost_per_mt', 2000):,}"],
        ["QC & Testing", f"{cfg.get('qc_cost_per_mt', 500):,}"],
        ["Overheads & Misc", f"{cfg.get('misc_cost_per_mt', 1000):,}"],
        ["TOTAL VARIABLE COST", f"{cfg.get('total_variable_cost_per_mt', 22000):,}"],
        ["Selling Price", f"{cfg.get('selling_price_per_mt', 35000):,}"],
        ["PROFIT PER MT", f"{cfg.get('profit_per_mt', 5000):,}"],
    ])

    # 15E. INSURANCE PLAN
    doc.add_heading("15.5 Insurance Plan", level=2)
    _add_table(doc, ["Insurance Type", "Coverage", "Premium (approx)"], [
        ["Fire + Allied Perils", f"Rs {cfg.get('investment_cr', 8):.2f} Cr (all fixed assets)", "0.15-0.25% of value"],
        ["Marine Transit", "Imported machinery in transit", "0.1-0.2%"],
        ["Public Liability", "Third party claims (hazardous substance)", "Rs 50,000-1,00,000/year"],
        ["Workmen Compensation", "All factory workers", "Rs 30,000-60,000/year"],
        ["Keyman Insurance", "Key promoter", "As per bank requirement"],
    ])

    # 15F. BANK BENCHMARK COMPARISON
    doc.add_heading("15.6 Bank Benchmark Comparison", level=2)
    _add_table(doc, ["Parameter", "Bank Minimum", "This Project", "Status"], [
        ["Promoter Contribution", "25%", f"{cfg.get('equity_ratio', 0.4)*100:.0f}%",
         "PASS" if cfg.get('equity_ratio', 0.4) >= 0.25 else "REVIEW"],
        ["DSCR (Average)", "1.50x", f"{cfg.get('dscr_yr3', 0):.2f}x",
         "PASS" if cfg.get('dscr_yr3', 0) >= 1.5 else "REVIEW"],
        ["IRR", "> 15%", f"{cfg.get('irr_pct', 0):.1f}%",
         "PASS" if cfg.get('irr_pct', 0) >= 15 else "REVIEW"],
        ["Break-Even", "< 60% capacity", f"{cfg.get('break_even_months', 0)} months",
         "PASS" if cfg.get('break_even_months', 0) <= 36 else "REVIEW"],
        ["Current Ratio", "> 1.33", f"{cfg.get('current_ratio', 1.5):.2f}",
         "PASS" if cfg.get('current_ratio', 1.5) >= 1.33 else "REVIEW"],
        ["D/E Ratio", "< 3.0", f"{cfg.get('debt_equity_ratio', 1.5):.2f}",
         "PASS" if cfg.get('debt_equity_ratio', 1.5) <= 3.0 else "REVIEW"],
    ])

    # 15G. PROJECT COST BREAKUP (BANK FORMAT)
    doc.add_heading("15.7 Project Cost Breakup (Standard Bank Format)", level=2)
    inv = cfg.get('investment_cr', 8)
    _add_table(doc, ["Component", "Rs Lac", "% of Total"], [
        ["Land & Site Development", f"{inv*100*0.12:.0f}", "12%"],
        ["Building & Civil Works", f"{inv*100*0.22:.0f}", "22%"],
        ["Plant & Machinery (Indigenous)", f"{inv*100*0.40:.0f}", "40%"],
        ["Miscellaneous Fixed Assets", f"{inv*100*0.04:.0f}", "4%"],
        ["Consultancy & Technical Fees", f"{inv*100*0.03:.0f}", "3%"],
        ["Pre-operative Expenses", f"{inv*100*0.04:.0f}", "4%"],
        ["Contingency (5%)", f"{inv*100*0.05:.0f}", "5%"],
        ["Working Capital Margin", f"{inv*100*0.10:.0f}", "10%"],
        ["TOTAL PROJECT COST", f"{inv*100:.0f}", "100%"],
    ])

    # 15H. CA CERTIFICATE TEMPLATE
    doc.add_heading("15.8 Chartered Accountant Certificate (Template)", level=2)
    doc.add_paragraph(
        "To Whom It May Concern\n\n"
        "I, [CA Name], (Membership No: [ICAI No], FRN: [Firm Reg No]), "
        "hereby certify that:\n\n"
        "1. The projected financial statements in this DPR are prepared on reasonable "
        "assumptions and accepted accounting principles.\n"
        "2. The net worth of the promoter(s) as on [date] is Rs [amount] as verified "
        "from audited books and supporting documents.\n"
        "3. The promoter(s) have no overdue loan/advance with any bank/financial institution.\n"
        "4. The DSCR, BEP, and IRR calculations are mathematically correct.\n"
        "5. The CMA data forms (I to VI) are prepared as per RBI guidelines.\n\n"
        "UDIN: [Auto-generated]\n"
        "Date: [Date]\n"
        "Place: [City]\n\n"
        "[CA Name]\n"
        "[Firm Name]\n"
        "[Membership No] | [FRN]\n"
        "[Seal & Signature]"
    )
    doc.add_page_break()

    _add_heading(doc, "16. CONSULTANT PROFILE")
    doc.add_paragraph(f"{company['trade_name']} — {company.get('tagline', '')}")
    doc.add_paragraph(f"Owner & Managing Director: {company['owner']}")
    doc.add_paragraph(f"Experience: {company.get('experience', '25 years')}")
    doc.add_paragraph(f"Headquarters: {company.get('hq', 'Vadodara, Gujarat')}")
    doc.add_paragraph(f"Contact: {company['phone']} | {company.get('email', '')}")
    doc.add_paragraph("")
    doc.add_paragraph("Key Credentials:")
    for cred in ["BSE-Listed Founder — Omnipotent Industries (1.2L MT, 11 JVs)",
                  "International Import Contracts — 2.4 Lakh MT/yr VG-30 (Iraq/USA)",
                  "10 Plants Built — 5 Product Types across 17 States",
                  "4,452 Industry Contacts — Contractors, Traders, Importers",
                  "Pride of India Award — Best Fast-Growing Business 2021"]:
        doc.add_paragraph(f"  - {cred}")

    # 17. Civil & Infrastructure Requirements
    _add_heading(doc, "17. CIVIL & INFRASTRUCTURE")
    try:
        from engines.plant_engineering import get_civil_specs
        civil = get_civil_specs(cfg)
        doc.add_paragraph(f"Process Hall: {civil['process_hall']['type']}, "
                          f"{civil['process_hall']['area_sqm']} sqm, "
                          f"{civil['process_hall']['length_m']}m x {civil['process_hall']['width_m']}m")
        doc.add_paragraph(f"Office: {civil['office_building']['type']}, {civil['office_building']['area_sqm']} sqm")
        doc.add_paragraph(f"Laboratory: {civil['laboratory']['area_sqm']} sqm, {civil['laboratory']['type']}")
        doc.add_paragraph(f"Control Room: {civil['control_room']['area_sqm']} sqm, blast-resistant {civil['control_room']['is_standard']}")
        doc.add_paragraph(f"Compound Wall: {civil['compound_wall']['total_length_m']}m perimeter, {civil['compound_wall']['height_m']}m height")
        doc.add_paragraph(f"Internal Roads: {civil['internal_roads']['width_m']}m wide, {civil['internal_roads']['pavement']}")
        doc.add_paragraph(f"Drainage: {civil['storm_drainage']['type']}, gradient {civil['storm_drainage']['gradient']}")
        doc.add_paragraph(f"Bund Wall: {civil['bund_wall']['height_m']}m x {civil['bund_wall']['thickness_m']}m, {civil['bund_wall']['material']}")
        doc.add_paragraph(f"Green Belt: {civil['green_belt']['width_m']}m wide, {civil['green_belt']['total_area_sqm']} sqm")
        doc.add_paragraph(f"Earthing: {civil['earthing_grid']['conductor']}, {civil['earthing_grid']['resistance']}")
    except Exception:
        doc.add_paragraph(f"Civil infrastructure based on {cfg.get('capacity_tpd', 20):.0f} TPD capacity. "
                          f"Plot: {cfg.get('plot_length_m', 120)}m x {cfg.get('plot_width_m', 80)}m")
    doc.add_page_break()

    # 18. Risk Assessment Matrix
    _add_heading(doc, "18. RISK ASSESSMENT & MITIGATION")
    _add_table(doc, ["Risk Category", "Risk Description", "Impact", "Probability", "Mitigation"],
        [
            ["Market", "Bitumen price volatility", "High", "Medium",
             "Diversified product mix (VG30+VG40+Char+Oil)"],
            ["Technology", "Process yield below design", "Medium", "Low",
             "CSIR-CRRI proven technology, pilot testing"],
            ["Financial", "Loan repayment stress", "High", "Low",
             f"DSCR {cfg.get('dscr_yr3', 0):.2f}x > 1.5x bank minimum"],
            ["Regulatory", "Environmental clearance delay", "Medium", "Medium",
             "Pre-submission consultation with SPCB"],
            ["Supply Chain", "Biomass availability seasonal", "Medium", "Medium",
             "90-day buffer stock, multiple feedstock sources"],
            ["Operational", "Skilled manpower shortage", "Low", "Medium",
             "Training program, operator certification"],
        ])
    doc.add_page_break()

    # 19. Covering Letter (for bank/govt submission)
    _add_heading(doc, "COVERING LETTER", level=0)
    doc.add_paragraph(f"Date: {now.strftime('%d %B %Y')}")
    doc.add_paragraph("")
    doc.add_paragraph("To,")
    doc.add_paragraph("The Branch Manager / Sanctioning Authority")
    doc.add_paragraph("[Bank Name / Government Authority]")
    doc.add_paragraph("[Address]")
    doc.add_paragraph("")
    doc.add_paragraph(f"Subject: Submission of Detailed Project Report — {project_name}")
    doc.add_paragraph("")
    doc.add_paragraph(f"Dear Sir/Madam,")
    doc.add_paragraph(
        f"We are pleased to submit the Detailed Project Report (DPR) for our proposed "
        f"{cfg.get('capacity_tpd', 20):.0f} MT/Day Bio-Modified Bitumen manufacturing plant at "
        f"{site_address}. The project requires a total investment of Rs {cfg.get('investment_cr', 8):.2f} Crore "
        f"with a term loan requirement of Rs {cfg.get('loan_cr', cfg.get('investment_cr', 8)*0.6):.2f} Crore."
    )
    doc.add_paragraph(
        f"The project utilizes CSIR-CRRI licensed pyrolysis technology for converting agricultural "
        f"waste into IS:73 compliant bio-modified bitumen. Key financial metrics: "
        f"ROI {cfg.get('roi_pct', 0):.1f}%, IRR {cfg.get('irr_pct', 0):.1f}%, "
        f"DSCR {cfg.get('dscr_yr3', 0):.2f}x, Break-even {cfg.get('break_even_months', 0)} months."
    )
    doc.add_paragraph(
        f"We request your kind consideration and sanction of the term loan facility. "
        f"All supporting documents including financial projections, machinery quotations, "
        f"and compliance checklist are enclosed herewith."
    )
    doc.add_paragraph("")
    doc.add_paragraph("Thanking you,")
    doc.add_paragraph(f"For {client_company or client_name}")
    doc.add_paragraph("")
    doc.add_paragraph(f"{client_name}")
    doc.add_paragraph("(Authorized Signatory)")
    doc.add_page_break()

    # 20. Annexures List
    _add_heading(doc, "20. LIST OF ANNEXURES")
    annexures = [
        "Annexure A: Promoter's KYC Documents (PAN, Aadhaar, Address Proof)",
        "Annexure B: Company Registration Certificate (ROC / MSME / Udyam)",
        "Annexure C: GST Registration Certificate",
        "Annexure D: Land Documents (Sale Deed / Lease Agreement / NOC)",
        "Annexure E: Machinery Quotations from OEM / Authorized Dealers",
        "Annexure F: Civil Construction Estimate from Licensed Contractor",
        "Annexure G: Consent to Establish (CTE) Application — State PCB",
        "Annexure H: Factory License Application — Factory Inspector",
        "Annexure I: Fire NOC Application — Fire Department",
        "Annexure J: PESO License Application (for petroleum storage)",
        "Annexure K: MSME / Udyam Registration",
        "Annexure L: Bank Account Statement (last 12 months)",
        "Annexure M: Income Tax Returns (last 3 years)",
        "Annexure N: CIBIL Report of Promoter(s)",
        "Annexure O: IS:73 Product Quality Test Certificate (from CSIR-CRRI lab)",
        "Annexure P: Site Photographs and Location Map",
        "Annexure Q: Equipment Supplier Catalogue / Technical Specifications",
        "Annexure R: Environmental Impact Assessment Summary",
        "Annexure S: Project Implementation Schedule (Gantt Chart)",
        "Annexure T: NHAI / MoRTH Specification Compliance Statement",
    ]
    for ann in annexures:
        doc.add_paragraph(ann, style='List Bullet')

    doc.add_paragraph("")
    doc.add_paragraph(
        "Note: All annexures to be submitted as per the specific requirements of the "
        "financing institution / government authority. Additional documents may be "
        "required based on state-specific regulations."
    )
    doc.add_page_break()

    # 21. Government Application Forms Required
    _add_heading(doc, "21. GOVERNMENT APPLICATIONS & FORMS REQUIRED")
    _add_table(doc, ["Application", "Authority", "Timeline", "Status"],
        [
            ["Consent to Establish (CTE)", "State Pollution Control Board", "30-60 days", "To Apply"],
            ["Consent to Operate (CTO)", "State Pollution Control Board", "15-30 days (after CTE)", "After Plant Ready"],
            ["Factory License", "Factory Inspector", "15-30 days", "To Apply"],
            ["Fire NOC", "Fire Department", "15-30 days", "To Apply"],
            ["PESO License", "Petroleum & Explosives Safety Org", "30-60 days", "To Apply"],
            ["Electrical Inspector Approval", "State Electrical Inspectorate", "15-30 days", "To Apply"],
            ["GST Registration", "GST Portal (gst.gov.in)", "7 days", "To Apply"],
            ["MSME / Udyam Registration", "Udyam Portal", "Instant", "To Apply"],
            ["GeM Portal Registration", "gem.gov.in", "3-5 days", "For Govt Supply"],
            ["CGTMSE Guarantee", "cgtmse.in", "Along with bank loan", "If Eligible"],
            [f"State Industrial Policy ({cfg.get('state', 'State')})", "State Industries Dept", "30-60 days", "To Apply"],
        ])
    doc.add_paragraph(
        f"All applications will be prepared and coordinated by {company['trade_name']} "
        f"as part of the consulting mandate."
    )
    doc.add_page_break()

    # Disclaimer
    doc.add_page_break()
    _add_heading(doc, "DISCLAIMER & CONFIDENTIALITY")
    doc.add_paragraph(
        "This Detailed Project Report is prepared by PPS Anantams Corporation Pvt Ltd based on "
        "current market data, verified engineering assumptions, and standard industry practices. "
        "All financial projections are estimates and subject to market conditions, regulatory changes, "
        "and operational efficiency. The client should conduct independent due diligence before "
        "making any investment decisions."
    )
    doc.add_paragraph(
        f"\nThis document is CONFIDENTIAL and prepared exclusively for {client_name}. "
        f"Unauthorized distribution or reproduction is prohibited."
    )
    doc.add_paragraph(f"\nPrepared by: {company['trade_name']}")
    doc.add_paragraph(f"{company['owner']} | {company['phone']} | {company['hq']}")
    doc.add_paragraph(f"GST: {company.get('gst', '')} | CIN: {company.get('cin', '')}")

    _footer(doc, company)
    return doc


# ═══════════════════════════════════════════════════════════════════════
# 2. BANK LOAN PROPOSAL — DOCX
# ═══════════════════════════════════════════════════════════════════════
def generate_bank_proposal_docx(cfg, company, customer=None):
    """Generate bank loan proposal DOCX."""
    doc = Document()
    now = datetime.now(IST)
    cust_name = customer.get("name", "The Branch Manager") if customer else "The Branch Manager"
    cust_company = customer.get("company", "") if customer else ""

    _add_heading(doc, "TERM LOAN PROPOSAL")
    doc.add_heading(f"Bio-Modified Bitumen Plant — {cfg.get('capacity_tpd', 20):.0f} MT/Day", level=2)
    doc.add_paragraph(f"Date: {now.strftime('%d %B %Y')}")
    doc.add_paragraph(f"To: {cust_name}")
    if cust_company:
        doc.add_paragraph(f"Company: {cust_company}")
    doc.add_paragraph(f"From: {company['trade_name']}")
    doc.add_paragraph("")

    _add_heading(doc, "1. LOAN REQUIREMENT")
    _add_table(doc, ["Parameter", "Details"], [
        ["Total Project Cost", f"Rs {cfg.get('investment_cr', 8):.2f} Crore"],
        ["Term Loan Required", f"Rs {cfg.get('loan_cr', 4.8):.2f} Crore"],
        ["Promoter Equity (40%)", f"Rs {cfg.get('equity_cr', 3.2):.2f} Crore"],
        ["Interest Rate", f"{cfg.get('interest_rate', 0.115)*100:.1f}% p.a."],
        ["Tenure", f"{cfg.get('emi_tenure_months', 84)} months ({cfg.get('emi_tenure_months', 84)//12} years)"],
        ["Monthly EMI", f"Rs {cfg.get('emi_lac_mth', 0):.2f} Lakhs"],
        ["DSCR (Year 3)", f"{cfg.get('dscr_yr3', 0):.2f}x (Bank minimum: 1.50x)"],
    ])

    _add_heading(doc, "2. PROJECT VIABILITY")
    _add_table(doc, ["Metric", "Value"], [
        ["Annual Revenue (Yr5)", f"Rs {cfg.get('revenue_yr5_lac', 0):.0f} Lakhs"],
        ["ROI", f"{cfg.get('roi_pct', 0):.1f}%"],
        ["IRR (Equity)", f"{cfg.get('irr_pct', 0):.1f}%"],
        ["Break-Even", f"{cfg.get('break_even_months', 0)} months"],
        ["Profit per MT", f"Rs {cfg.get('profit_per_mt', 0):,.0f}"],
    ])

    _add_heading(doc, "3. SECURITY OFFERED")
    doc.add_paragraph("• Primary: Hypothecation of Plant & Machinery, Current Assets")
    doc.add_paragraph("• Collateral: Immovable property (Land + Building)")
    doc.add_paragraph("• Personal Guarantee of Promoter Directors")
    doc.add_paragraph(f"• Promoter: {company['owner']} | {company['experience']}")

    _footer(doc, company)
    return doc


# ═══════════════════════════════════════════════════════════════════════
# 3. INVESTOR PITCH — PPTX
# ═══════════════════════════════════════════════════════════════════════
def generate_investor_pptx(cfg, company):
    """Generate investor pitch deck PPTX."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor as PptxRGB

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    def add_slide(title_text, content_text=""):
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
        slide.shapes.title.text = title_text
        if content_text and len(slide.placeholders) > 1:
            slide.placeholders[1].text = content_text
        return slide

    # Slide 1: Cover
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = f"{company['trade_name']}"
    s1.placeholders[1].text = (
        f"Bio-Modified Bitumen Plant — {cfg.get('capacity_tpd', 20):.0f} MT/Day\n"
        f"Investment: Rs {cfg.get('investment_cr', 8):.2f} Crore | IRR: {cfg.get('irr_pct', 0):.1f}%\n"
        f"{company['owner']} | {company['phone']}"
    )

    # Slide 2: Opportunity
    add_slide("The Opportunity", (
        f"• India's road network: 6.4 million km\n"
        f"• Govt allocation: Rs 2.7 lakh crore for highways (FY26)\n"
        f"• Bitumen demand growing at 8-10% CAGR\n"
        f"• Green mandate from NHAI for sustainable materials\n"
        f"• Bio-Bitumen: 15-30% cheaper than conventional VG30"
    ))

    # Slide 3: Investment Summary
    add_slide("Investment Summary", (
        f"• Total Project Cost: Rs {cfg.get('investment_cr', 8):.2f} Crore\n"
        f"• Equity Required: Rs {cfg.get('equity_cr', 3.2):.2f} Crore\n"
        f"• Bank Loan: Rs {cfg.get('loan_cr', 4.8):.2f} Crore\n"
        f"• Annual Revenue (Yr5): Rs {cfg.get('revenue_yr5_lac', 0):.0f} Lakhs\n"
        f"• ROI: {cfg.get('roi_pct', 0):.1f}% | IRR: {cfg.get('irr_pct', 0):.1f}%\n"
        f"• Break-Even: {cfg.get('break_even_months', 0)} months\n"
        f"• DSCR: {cfg.get('dscr_yr3', 0):.2f}x"
    ))

    # Slide 4: Technology
    add_slide("Technology: Biomass Pyrolysis", (
        f"• Agro-waste (rice straw, bagasse) → Pyrolysis at 400-600°C\n"
        f"• Bio-Oil (40% yield) → Blended with bitumen\n"
        f"• Biochar (30%) → Soil amendment / Carbon black\n"
        f"• Syngas (25%) → Captive fuel (saves Rs 50+ Lakh/year)\n"
        f"• CSIR-CRRI validated technology\n"
        f"• BIS & MoRTH 2026 compliant"
    ))

    # Slide 5: Plant Specs
    add_slide(f"Plant: {cfg.get('capacity_tpd', 20):.0f} MT/Day", (
        f"• Staff: {cfg.get('staff', 20)} persons\n"
        f"• Power: {cfg.get('power_kw', 100):.0f} kW\n"
        f"• Oil Output: {cfg.get('oil_ltr_day', 8000):.0f} Litres/day\n"
        f"• Char Output: {cfg.get('char_kg_day', 6000):.0f} Kg/day\n"
        f"• Land: {int(cfg.get('capacity_tpd', 20) * 520):,} sq ft\n"
        f"• Location: {cfg.get('location', 'TBD')}, {cfg.get('state', '')}"
    ))

    # Slide 6: Financial Projection
    if cfg.get("roi_timeline"):
        content = "Year | Revenue | EBITDA | PAT | DSCR\n"
        for yr in cfg.get("roi_timeline", 0)[:5]:
            content += f"Yr{yr['Year']} | Rs {yr['Revenue (Lac)']:.0f}L | Rs {yr['EBITDA (Lac)']:.0f}L | Rs {yr['PAT (Lac)']:.0f}L | {yr['DSCR']:.2f}x\n"
        add_slide("7-Year Financial Projection", content)

    # Slide 7: Exit Strategy
    add_slide("Exit Strategy", (
        f"• Year 3-5: Promoter buyback at 2x equity\n"
        f"• Year 5-7: Strategic sale to infrastructure company\n"
        f"• IPO route: After 3+ years of profitable operations\n"
        f"• Expected equity IRR: {cfg.get('irr_pct', 0):.1f}%\n"
        f"• Payback period: {cfg.get('break_even_months', 0)} months"
    ))

    # Slide 8: Contact
    add_slide("Contact Us", (
        f"{company['name']}\n"
        f"{company['owner']}\n"
        f"Phone: {company['phone']}\n"
        f"Email: {company.get('email', '')}\n"
        f"Location: {company['hq']}\n"
        f"GST: {company['gst']}\n\n"
        f"CONFIDENTIAL — For Private Circulation Only"
    ))

    return prs


# ═══════════════════════════════════════════════════════════════════════
# 4. FINANCIAL MODEL — XLSX
# ═══════════════════════════════════════════════════════════════════════
def generate_financial_xlsx(cfg, company):
    """Generate complete financial model Excel workbook."""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, numbers

    wb = Workbook()
    header_font = Font(bold=True, color="FFFFFF", size=11)
    header_fill = PatternFill(start_color="003366", end_color="003366", fill_type="solid")
    num_fmt = '#,##0.00'

    def write_header(ws, row, cols):
        for i, c in enumerate(cols):
            cell = ws.cell(row=row, column=i + 1, value=c)
            cell.font = header_font
            cell.fill = header_fill

    # Sheet 1: Summary
    ws1 = wb.active
    ws1.title = "Summary"
    summary = [
        [f"Bio Bitumen Financial Model — {cfg.get('capacity_tpd', 20):.0f} MT/Day", ""],
        ["Date", datetime.now(IST).strftime("%d %B %Y")],
        ["", ""],
        ["INVESTMENT", ""],
        ["Total Project Cost", f"Rs {cfg.get('investment_cr', 8):.2f} Crore"],
        ["Bank Loan (60%)", f"Rs {cfg.get('loan_cr', 4.8):.2f} Crore"],
        ["Equity (40%)", f"Rs {cfg.get('equity_cr', 3.2):.2f} Crore"],
        ["Monthly EMI", f"Rs {cfg.get('emi_lac_mth', 0):.2f} Lakhs"],
        ["", ""],
        ["REVENUE", ""],
        ["Selling Price/MT", f"Rs {cfg.get('selling_price_per_mt', 35000):,}"],
        ["Variable Cost/MT", f"Rs {cfg.get('total_variable_cost_per_mt', 22000):,}"],
        ["Profit/MT", f"Rs {cfg.get('profit_per_mt', 0):,}"],
        ["Revenue Yr5", f"Rs {cfg.get('revenue_yr5_lac', 0):.0f} Lakhs"],
        ["Monthly Profit", f"Rs {cfg.get('monthly_profit_lac', 0):.1f} Lakhs"],
        ["", ""],
        ["RETURNS", ""],
        ["ROI", f"{cfg.get('roi_pct', 0):.1f}%"],
        ["IRR", f"{cfg.get('irr_pct', 0):.1f}%"],
        ["DSCR Yr3", f"{cfg.get('dscr_yr3', 0):.2f}x"],
        ["Break-Even", f"{cfg.get('break_even_months', 0)} months"],
    ]
    for i, (a, b) in enumerate(summary):
        ws1.cell(row=i + 1, column=1, value=a).font = Font(bold=True) if not b else Font()
        ws1.cell(row=i + 1, column=2, value=b)

    # Sheet 2: 7-Year P&L
    ws2 = wb.create_sheet("P&L_7Year")
    if cfg.get("roi_timeline"):
        headers = list(cfg.get("roi_timeline", 0)[0].keys())
        write_header(ws2, 1, headers)
        for i, yr in enumerate(cfg.get("roi_timeline", 0)):
            for j, h in enumerate(headers):
                ws2.cell(row=i + 2, column=j + 1, value=yr[h])

    # Sheet 3: Cost Structure
    ws3 = wb.create_sheet("Cost_Structure")
    write_header(ws3, 1, ["Cost Item", "Rs per MT"])
    costs = [
        ("Raw Material", cfg.get("raw_material_cost_per_mt", 0)),
        ("Power & Fuel", cfg.get("power_cost_per_mt", 0)),
        ("Labour", cfg.get("labour_cost_per_mt", 0)),
        ("Chemicals", cfg.get("chemical_cost_per_mt", 0)),
        ("Packaging", cfg.get("packaging_cost_per_mt", 0)),
        ("Transport", cfg.get("transport_cost_per_mt", 0)),
        ("QC & Testing", cfg.get("qc_cost_per_mt", 0)),
        ("Miscellaneous", cfg.get("misc_cost_per_mt", 0)),
        ("TOTAL", cfg.get("total_variable_cost_per_mt", 0)),
    ]
    for i, (item, val) in enumerate(costs):
        ws3.cell(row=i + 2, column=1, value=item)
        ws3.cell(row=i + 2, column=2, value=val)

    # Sheet 4: Sensitivity
    ws4 = wb.create_sheet("Sensitivity")
    if cfg.get("sensitivity_matrix"):
        write_header(ws4, 1, ["Cost \\ Price", "Low Price", "Base Price", "High Price"])
        labels = ["Low Cost", "Base Cost", "High Cost"]
        for i, (label, row) in enumerate(zip(labels, cfg.get("sensitivity_matrix", 0))):
            ws4.cell(row=i + 2, column=1, value=label).font = Font(bold=True)
            for j, val in enumerate(row):
                ws4.cell(row=i + 2, column=j + 2, value=val)

    # Sheet 5: Assumptions
    ws5 = wb.create_sheet("Assumptions")
    write_header(ws5, 1, ["Parameter", "Value", "Source"])
    assumptions = [
        ("Capacity", f"{cfg.get('capacity_tpd', 20):.0f} MT/Day", "Client requirement"),
        ("Working Days", str(cfg.get("working_days", 0)), "Industry standard"),
        ("Selling Price", f"Rs {cfg.get('selling_price_per_mt', 35000):,}/MT", "PetroBazaar Mar 2026"),
        ("Biomass Cost", f"Rs {cfg.get('raw_material_cost_per_mt', 8000):,}/MT output", "IndiaMART verified"),
        ("Interest Rate", f"{cfg.get('interest_rate', 0.115)*100:.1f}%", "SBI MCLR + 200 bps"),
        ("Tax Rate", f"{cfg.get('tax_rate', 0.25)*100:.0f}%", "Section 115BAB"),
        ("Depreciation", f"{cfg.get('depreciation_rate', 0.10)*100:.0f}% SLM", "Companies Act 2013"),
        ("Oil Yield", "40% of biomass", "CSIR-CRRI research"),
        ("Char Yield", "30% of biomass", "CSIR-CRRI research"),
    ]
    for i, (param, val, src) in enumerate(assumptions):
        ws5.cell(row=i + 2, column=1, value=param)
        ws5.cell(row=i + 2, column=2, value=val)
        ws5.cell(row=i + 2, column=3, value=src)

    return wb


# ═══════════════════════════════════════════════════════════════════════
# 5. MASTER GENERATOR — All documents at once
# ═══════════════════════════════════════════════════════════════════════
def generate_all_documents(cfg, company, customer=None, output_dir=None):
    """
    Generate ALL documents for current config.
    Returns dict of {filename: bytes_or_path}
    """
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)

    cap_label = f"{cfg.get('capacity_tpd', 20):.0f}MT"
    cust_label = customer["name"].replace(" ", "_") if customer else "General"
    results = {}

    # 1. DPR DOCX
    dpr_doc = generate_dpr_docx(cfg, company)
    fname = f"DPR_{cap_label}_{cust_label}.docx"
    if output_dir:
        path = os.path.join(output_dir, fname)
        dpr_doc.save(path)
        results[fname] = path
    else:
        buf = io.BytesIO()
        dpr_doc.save(buf)
        results[fname] = buf.getvalue()

    # 2. Bank Proposal DOCX
    bank_doc = generate_bank_proposal_docx(cfg, company, customer)
    fname = f"Bank_Proposal_{cap_label}_{cust_label}.docx"
    if output_dir:
        path = os.path.join(output_dir, fname)
        bank_doc.save(path)
        results[fname] = path
    else:
        buf = io.BytesIO()
        bank_doc.save(buf)
        results[fname] = buf.getvalue()

    # 3. Investor PPTX
    pptx = generate_investor_pptx(cfg, company)
    fname = f"Investor_Pitch_{cap_label}_{cust_label}.pptx"
    if output_dir:
        path = os.path.join(output_dir, fname)
        pptx.save(path)
        results[fname] = path
    else:
        buf = io.BytesIO()
        pptx.save(buf)
        results[fname] = buf.getvalue()

    # 4. Financial Model XLSX
    xlsx = generate_financial_xlsx(cfg, company)
    fname = f"Financial_Model_{cap_label}_{cust_label}.xlsx"
    if output_dir:
        path = os.path.join(output_dir, fname)
        xlsx.save(path)
        results[fname] = path
    else:
        buf = io.BytesIO()
        xlsx.save(buf)
        results[fname] = buf.getvalue()

    return results
